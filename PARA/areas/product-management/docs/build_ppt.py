from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Paleta Interbank ──────────────────────────────────────────────────────────
AZUL_OSC  = RGBColor(0x00, 0x20, 0x60)
AZUL_MED  = RGBColor(0x00, 0x39, 0xA6)
VERDE_PRI = RGBColor(0x05, 0xBE, 0x50)
VERDE_OSC = RGBColor(0x00, 0x8C, 0x37)
BLANCO    = RGBColor(0xFF, 0xFF, 0xFF)
NEGRO     = RGBColor(0x00, 0x00, 0x00)
GRIS_SUB  = RGBColor(0x60, 0x60, 0x60)
ROJO_ALERT= RGBColor(0xC0, 0x00, 0x00)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
blank = prs.slide_layouts[6]  # blank layout

def add_slide():
    return prs.slides.add_slide(blank)

def rect(slide, x, y, w, h, fill=None, line=None):
    from pptx.util import Emu
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    shp = slide.shapes.add_shape(1, x, y, w, h)  # MSO_SHAPE_TYPE.RECTANGLE=1
    shp.line.fill.background()
    if fill:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    else:
        shp.fill.background()
    if line:
        shp.line.color.rgb = line
        shp.line.width = Pt(1)
    else:
        shp.line.fill.background()
    return shp

def txbox(slide, text, x, y, w, h, size=16, bold=False, color=NEGRO,
          align=PP_ALIGN.LEFT, font="Poppins", wrap=True):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tb.word_wrap = wrap
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb

def card(slide, x, y, w, h, bg, title, title_size, title_color,
         body=None, body_size=14, body_color=NEGRO):
    rect(slide, x, y, w, h, fill=bg)
    margin = Inches(0.2)
    txbox(slide, title, x+margin, y+margin, w-2*margin, Inches(0.6),
          size=title_size, bold=True, color=title_color, align=PP_ALIGN.CENTER)
    if body:
        txbox(slide, body, x+margin, y+Inches(0.7), w-2*margin, h-Inches(0.9),
              size=body_size, color=body_color, align=PP_ALIGN.CENTER)

# ── SLIDE 1: PORTADA ─────────────────────────────────────────────────────────
sl1 = add_slide()
rect(sl1, 0, 0, W, H, fill=AZUL_OSC)

# Círculo verde decorativo arriba derecha
circ = sl1.shapes.add_shape(9, Inches(10.5), Inches(-1), Inches(4), Inches(4))  # oval=9
circ.fill.solid(); circ.fill.fore_color.rgb = AZUL_MED
circ.line.fill.background()

circ2 = sl1.shapes.add_shape(9, Inches(11.5), Inches(5), Inches(3.5), Inches(3.5))
circ2.fill.solid(); circ2.fill.fore_color.rgb = VERDE_OSC
circ2.line.fill.background()

txbox(sl1, "SEGUIMIENTO DE COMPROMISOS", Inches(0.8), Inches(1.8), Inches(9), Inches(1.0),
      size=14, bold=False, color=VERDE_PRI, align=PP_ALIGN.LEFT)
txbox(sl1, "Depósitos · Abril 2026", Inches(0.8), Inches(2.5), Inches(10), Inches(1.6),
      size=44, bold=True, color=BLANCO, align=PP_ALIGN.LEFT)
txbox(sl1, "Análisis de incumplimiento por Banca, Zonal y Ejecutivo",
      Inches(0.8), Inches(4.1), Inches(9), Inches(0.8),
      size=18, bold=False, color=RGBColor(0xCC, 0xDD, 0xFF), align=PP_ALIGN.LEFT)
txbox(sl1, "Interbank · Product Management", Inches(0.8), Inches(6.5), Inches(6), Inches(0.5),
      size=11, bold=False, color=RGBColor(0x99, 0xAA, 0xCC), align=PP_ALIGN.LEFT)

# ── SLIDE 2: RESUMEN EJECUTIVO ───────────────────────────────────────────────
sl2 = add_slide()
rect(sl2, 0, 0, W, Inches(1.3), fill=AZUL_OSC)
txbox(sl2, "Resumen Ejecutivo", Inches(0.5), Inches(0.3), Inches(10), Inches(0.7),
      size=28, bold=True, color=BLANCO, align=PP_ALIGN.LEFT)
txbox(sl2, "Todos los clientes monitoreados incumplen su compromiso · Período: Abril 2026",
      Inches(0.5), Inches(1.5), Inches(12.3), Inches(0.5),
      size=13, bold=False, color=GRIS_SUB)

# 3 KPI cards
kpis = [
    ("103", "Clientes\nincumpliendo", AZUL_OSC),
    ("38.6%", "Cumplimiento\npromedio", VERDE_OSC),
    ("S/ 2.1M", "Ganancia\npotencial perdida", ROJO_ALERT),
]
cw = Inches(3.8)
gap = Inches(0.3)
start_x = Inches(0.5)
for i, (num, lbl, col) in enumerate(kpis):
    cx = start_x + i*(cw+gap)
    rect(sl2, cx, Inches(2.1), cw, Inches(3.8), fill=col)
    txbox(sl2, num, cx, Inches(2.5), cw, Inches(1.8),
          size=64, bold=True, color=BLANCO, align=PP_ALIGN.CENTER)
    txbox(sl2, lbl, cx, Inches(4.3), cw, Inches(1.0),
          size=16, bold=False, color=BLANCO, align=PP_ALIGN.CENTER)

txbox(sl2, "Fuente: Excel reunion.xlsx · Procesado 2026-06-15",
      Inches(0.5), Inches(6.9), Inches(8), Inches(0.4),
      size=9, color=GRIS_SUB)

# ── SLIDE 3: POR BANCA ───────────────────────────────────────────────────────
sl3 = add_slide()
rect(sl3, 0, 0, W, Inches(1.3), fill=AZUL_OSC)
txbox(sl3, "Incumplimiento por Banca", Inches(0.5), Inches(0.3), Inches(10), Inches(0.7),
      size=28, bold=True, color=BLANCO, align=PP_ALIGN.LEFT)

# Tabla visual: 3 columnas (una por banca) con métricas apiladas
bancas = [
    ("BC", "Banca Corporativa", 53, "24.6%", "S/ 338M", "S/ 1,140M", "S/ 1.99M", AZUL_OSC),
    ("BEL", "Banca Empresa Líder", 19, "51.8%", "S/ 39M",  "S/ 72M",   "S/ 49.8K", AZUL_MED),
    ("BEP", "Banca Empresa Plus", 14, "69.8%", "S/ 433M", "S/ 528M",  "S/ 70.6K", VERDE_OSC),
]

col_w = Inches(3.9)
col_gap = Inches(0.28)
bx = Inches(0.4)

for i, (code, name, clientes, cumpl, saldo, comprom, extra, color) in enumerate(bancas):
    cx = bx + i*(col_w+col_gap)
    # Header card
    rect(sl3, cx, Inches(1.5), col_w, Inches(1.0), fill=color)
    txbox(sl3, code, cx, Inches(1.52), col_w, Inches(0.55),
          size=30, bold=True, color=BLANCO, align=PP_ALIGN.CENTER)
    txbox(sl3, name, cx, Inches(2.0), col_w, Inches(0.38),
          size=10, color=BLANCO, align=PP_ALIGN.CENTER)
    # Métricas
    metrics = [
        ("Clientes", str(clientes)),
        ("Cumplimiento prom.", cumpl),
        ("Saldo promedio total", saldo),
        ("Compromiso total", comprom),
        ("Extra perdido", extra),
    ]
    for j, (lbl, val) in enumerate(metrics):
        my = Inches(2.65) + j*Inches(0.72)
        bg = RGBColor(0xF0, 0xF4, 0xFF) if j%2==0 else BLANCO
        rect(sl3, cx, my, col_w, Inches(0.68), fill=bg)
        txbox(sl3, lbl, cx+Inches(0.1), my+Inches(0.04), col_w*0.55, Inches(0.35),
              size=11, color=GRIS_SUB)
        val_color = ROJO_ALERT if lbl == "Extra perdido" else AZUL_OSC
        txbox(sl3, val, cx+col_w*0.55, my+Inches(0.04), col_w*0.43, Inches(0.35),
              size=12, bold=True, color=val_color, align=PP_ALIGN.RIGHT)

txbox(sl3, "BC concentra el 95% de la ganancia potencial no realizada",
      Inches(0.4), Inches(6.8), Inches(12), Inches(0.4),
      size=11, bold=True, color=ROJO_ALERT)

# ── SLIDE 4: DETALLE BC ──────────────────────────────────────────────────────
sl4 = add_slide()
rect(sl4, 0, 0, W, Inches(1.3), fill=AZUL_OSC)
txbox(sl4, "Detalle BC — Banca Corporativa", Inches(0.5), Inches(0.3), Inches(10), Inches(0.7),
      size=28, bold=True, color=BLANCO)

# Subtítulo zonales
txbox(sl4, "Por Zonal", Inches(0.5), Inches(1.4), Inches(5), Inches(0.4),
      size=14, bold=True, color=AZUL_OSC)

zonales_bc = [
    ("BC GRUPO 1", 5, "35.2%", "S/ 17M",  "S/ 19.6K"),
    ("BC GRUPO 2", 48, "23.5%", "S/ 322M", "S/ 1.97M"),
]
headers = ["Zonal", "Clientes", "Cumpl.", "Saldo Total", "Extra Perdido"]
col_ws = [Inches(3.0), Inches(1.0), Inches(1.2), Inches(2.0), Inches(2.0)]
xs = [Inches(0.4)]
for cw2 in col_ws[:-1]:
    xs.append(xs[-1]+cw2+Inches(0.05))

# Header row
for j, (hdr, cw2, x2) in enumerate(zip(headers, col_ws, xs)):
    rect(sl4, x2, Inches(1.85), cw2, Inches(0.38), fill=AZUL_MED)
    txbox(sl4, hdr, x2+Inches(0.05), Inches(1.87), cw2-Inches(0.1), Inches(0.35),
          size=11, bold=True, color=BLANCO)

for i, (zonal, cli, cumpl, saldo, extra) in enumerate(zonales_bc):
    by = Inches(2.27) + i*Inches(0.45)
    bg = RGBColor(0xF0, 0xF4, 0xFF) if i%2==0 else BLANCO
    for cw2, x2 in zip(col_ws, xs):
        rect(sl4, x2, by, cw2, Inches(0.42), fill=bg)
    vals = [zonal, str(cli), cumpl, saldo, extra]
    for j, (v, cw2, x2) in enumerate(zip(vals, col_ws, xs)):
        vc = ROJO_ALERT if j==4 else NEGRO
        txbox(sl4, v, x2+Inches(0.05), by+Inches(0.05), cw2-Inches(0.1), Inches(0.35),
              size=12, color=vc, bold=(j==4))

# Top ejecutivos BC
txbox(sl4, "Top Ejecutivos por Ganancia Perdida", Inches(0.5), Inches(3.4), Inches(12), Inches(0.4),
      size=14, bold=True, color=AZUL_OSC)

ejecs_bc = [
    ("RODRÍGUEZ SANGURIMA M. L.", "BC GRUPO 2", 29, "17.4%", "S/ 1.93M"),
    ("ROMERO MASÍAS M.",           "BC GRUPO 2", 16, "24.9%", "S/ 35.9K"),
    ("DUHART NARANJO L. M.",       "BC GRUPO 1",  2, "20.7%", "S/ 15.7K"),
    ("MORA HARTMANN R. M.",        "BC GRUPO 2",  3, "42.5%", "S/ 3.9K"),
    ("ESCUDERO BERNARDO G.",       "BC GRUPO 2",  3, "77.7%", "S/ 2.5K"),
]
ehdrs = ["Ejecutivo", "Zonal", "Clientes", "Cumpl.", "Extra Perdido"]
ecws  = [Inches(4.2), Inches(2.2), Inches(1.0), Inches(1.2), Inches(1.8)]
exs   = [Inches(0.4)]
for cw2 in ecws[:-1]:
    exs.append(exs[-1]+cw2+Inches(0.05))

for j, (hdr, cw2, x2) in enumerate(zip(ehdrs, ecws, exs)):
    rect(sl4, x2, Inches(3.85), cw2, Inches(0.35), fill=VERDE_OSC)
    txbox(sl4, hdr, x2+Inches(0.05), Inches(3.87), cw2-Inches(0.1), Inches(0.32),
          size=10, bold=True, color=BLANCO)

for i, (ejec, zonal, cli, cumpl, extra) in enumerate(ejecs_bc):
    ey = Inches(4.23) + i*Inches(0.4)
    bg = RGBColor(0xF5, 0xFF, 0xF5) if i%2==0 else BLANCO
    for cw2, x2 in zip(ecws, exs):
        rect(sl4, x2, ey, cw2, Inches(0.38), fill=bg)
    vals = [ejec, zonal, str(cli), cumpl, extra]
    for j, (v, cw2, x2) in enumerate(zip(vals, ecws, exs)):
        vc = ROJO_ALERT if j==4 else NEGRO
        txbox(sl4, v, x2+Inches(0.05), ey+Inches(0.04), cw2-Inches(0.1), Inches(0.32),
              size=11, color=vc, bold=(j==4))

# ── SLIDE 5: DETALLE BEL ─────────────────────────────────────────────────────
sl5 = add_slide()
rect(sl5, 0, 0, W, Inches(1.3), fill=AZUL_MED)
txbox(sl5, "Detalle BEL — Banca Empresa Líder", Inches(0.5), Inches(0.3), Inches(10), Inches(0.7),
      size=28, bold=True, color=BLANCO)

txbox(sl5, "Por Zonal", Inches(0.5), Inches(1.4), Inches(5), Inches(0.4),
      size=14, bold=True, color=AZUL_OSC)

zonales_bel = [
    ("BEL ZONAL 1", 7, "45.5%", "S/ 17M", "S/ 28.3K"),
    ("BEL ZONAL 2", 3, "39.4%", "S/ 14M", "S/ 18.1K"),
    ("BEL ZONAL 3", 9, "61.6%", "S/ 7M",  "S/ 3.4K"),
]

for j, (hdr, cw2, x2) in enumerate(zip(headers, col_ws, xs)):
    rect(sl5, x2, Inches(1.85), cw2, Inches(0.38), fill=AZUL_MED)
    txbox(sl5, hdr, x2+Inches(0.05), Inches(1.87), cw2-Inches(0.1), Inches(0.35),
          size=11, bold=True, color=BLANCO)

for i, (zonal, cli, cumpl, saldo, extra) in enumerate(zonales_bel):
    by = Inches(2.27) + i*Inches(0.45)
    bg = RGBColor(0xF0, 0xF4, 0xFF) if i%2==0 else BLANCO
    for cw2, x2 in zip(col_ws, xs):
        rect(sl5, x2, by, cw2, Inches(0.42), fill=bg)
    vals = [zonal, str(cli), cumpl, saldo, extra]
    for j, (v, cw2, x2) in enumerate(zip(vals, col_ws, xs)):
        vc = ROJO_ALERT if j==4 else NEGRO
        txbox(sl5, v, x2+Inches(0.05), by+Inches(0.05), cw2-Inches(0.1), Inches(0.35),
              size=12, color=vc, bold=(j==4))

txbox(sl5, "Top Ejecutivos por Ganancia Perdida", Inches(0.5), Inches(3.5), Inches(12), Inches(0.4),
      size=14, bold=True, color=AZUL_OSC)

ejecs_bel = [
    ("CABRERA GALÁN DE NOLASCO Z. Y.", "BEL ZONAL 1", 2, "22.5%", "S/ 11.9K"),
    ("CHANG PILLIHUARI D. Z.",          "BEL ZONAL 2", 2, "49.5%", "S/ 9.5K"),
    ("BOLAÑOS DÍAZ R. F.",              "BEL ZONAL 2", 2, "49.0%", "S/ 8.6K"),
    ("MEDRANO LUNA A. A.",              "BEL ZONAL 1", 1, "19.7%", "S/ 8.3K"),
    ("ZEGARRA LLOSA A.",                "BEL ZONAL 3", 8, "59.5%", "S/ 1.9K"),
]

for j, (hdr, cw2, x2) in enumerate(zip(ehdrs, ecws, exs)):
    rect(sl5, x2, Inches(3.95), cw2, Inches(0.35), fill=AZUL_MED)
    txbox(sl5, hdr, x2+Inches(0.05), Inches(3.97), cw2-Inches(0.1), Inches(0.32),
          size=10, bold=True, color=BLANCO)

for i, (ejec, zonal, cli, cumpl, extra) in enumerate(ejecs_bel):
    ey = Inches(4.33) + i*Inches(0.4)
    bg = RGBColor(0xF5, 0xFF, 0xF5) if i%2==0 else BLANCO
    for cw2, x2 in zip(ecws, exs):
        rect(sl5, x2, ey, cw2, Inches(0.38), fill=bg)
    vals = [ejec, zonal, str(cli), cumpl, extra]
    for j, (v, cw2, x2) in enumerate(zip(vals, ecws, exs)):
        vc = ROJO_ALERT if j==4 else NEGRO
        txbox(sl5, v, x2+Inches(0.05), ey+Inches(0.04), cw2-Inches(0.1), Inches(0.32),
              size=11, color=vc, bold=(j==4))

# ── SLIDE 6: DETALLE BEP ─────────────────────────────────────────────────────
sl6 = add_slide()
rect(sl6, 0, 0, W, Inches(1.3), fill=VERDE_OSC)
txbox(sl6, "Detalle BEP — Banca Empresa Plus", Inches(0.5), Inches(0.3), Inches(10), Inches(0.7),
      size=28, bold=True, color=BLANCO)

txbox(sl6, "Por Zonal", Inches(0.5), Inches(1.4), Inches(5), Inches(0.4),
      size=14, bold=True, color=AZUL_OSC)

zonales_bep = [
    ("BEP ZONAL 1", 6, "82.9%", "S/ 362M", "S/ 19.0K"),
    ("BEP ZONAL 3", 8, "61.0%", "S/ 71M",  "S/ 51.6K"),
]

for j, (hdr, cw2, x2) in enumerate(zip(headers, col_ws, xs)):
    rect(sl6, x2, Inches(1.85), cw2, Inches(0.38), fill=VERDE_OSC)
    txbox(sl6, hdr, x2+Inches(0.05), Inches(1.87), cw2-Inches(0.1), Inches(0.35),
          size=11, bold=True, color=BLANCO)

for i, (zonal, cli, cumpl, saldo, extra) in enumerate(zonales_bep):
    by = Inches(2.27) + i*Inches(0.45)
    bg = RGBColor(0xF0, 0xFF, 0xF0) if i%2==0 else BLANCO
    for cw2, x2 in zip(col_ws, xs):
        rect(sl6, x2, by, cw2, Inches(0.42), fill=bg)
    vals = [zonal, str(cli), cumpl, saldo, extra]
    for j, (v, cw2, x2) in enumerate(zip(vals, col_ws, xs)):
        vc = ROJO_ALERT if j==4 else NEGRO
        txbox(sl6, v, x2+Inches(0.05), by+Inches(0.05), cw2-Inches(0.1), Inches(0.35),
              size=12, color=vc, bold=(j==4))

txbox(sl6, "Top Ejecutivos por Ganancia Perdida", Inches(0.5), Inches(3.4), Inches(12), Inches(0.4),
      size=14, bold=True, color=AZUL_OSC)

ejecs_bep = [
    ("GESTIÓN RAMÍREZ M.",          "BEP ZONAL 3", 6, "54.7%", "S/ 51.1K"),
    ("CALDERÓN CÓRDOVA K.",         "BEP ZONAL 3", 2, "54.5%", "S/ 17.6K"),
    ("ARANCIBIA BIANCO R. M.",      "BEP ZONAL 1", 4, "92.4%", "S/ 1.4K"),
]

for j, (hdr, cw2, x2) in enumerate(zip(ehdrs, ecws, exs)):
    rect(sl6, x2, Inches(3.85), cw2, Inches(0.35), fill=VERDE_OSC)
    txbox(sl6, hdr, x2+Inches(0.05), Inches(3.87), cw2-Inches(0.1), Inches(0.32),
          size=10, bold=True, color=BLANCO)

for i, (ejec, zonal, cli, cumpl, extra) in enumerate(ejecs_bep):
    ey = Inches(4.23) + i*Inches(0.4)
    bg = RGBColor(0xF0, 0xFF, 0xF0) if i%2==0 else BLANCO
    for cw2, x2 in zip(ecws, exs):
        rect(sl6, x2, ey, cw2, Inches(0.38), fill=bg)
    vals = [ejec, zonal, str(cli), cumpl, extra]
    for j, (v, cw2, x2) in enumerate(zip(vals, ecws, exs)):
        vc = ROJO_ALERT if j==4 else NEGRO
        txbox(sl6, v, x2+Inches(0.05), ey+Inches(0.04), cw2-Inches(0.1), Inches(0.32),
              size=11, color=vc, bold=(j==4))

# Nota BEP ZONAL 2 sin data
txbox(sl6, "Nota: BEP ZONAL 2 sin registros en el período.",
      Inches(0.4), Inches(6.8), Inches(10), Inches(0.4),
      size=10, color=GRIS_SUB)

# ── SLIDE 7: RANKING EJECUTIVOS ───────────────────────────────────────────────
sl7 = add_slide()
rect(sl7, 0, 0, W, Inches(1.3), fill=AZUL_OSC)
txbox(sl7, "Ranking de Ejecutivos — Ganancia Potencial No Realizada",
      Inches(0.5), Inches(0.25), Inches(12), Inches(0.85),
      size=24, bold=True, color=BLANCO)

all_ejecs = [
    (1, "RODRÍGUEZ SANGURIMA M. L.", "BC", "BC GRUPO 2", 29, 17.4, 1933255),
    (2, "GESTIÓN RAMÍREZ M.",         "BEP","BEP ZONAL 3", 6, 54.7,  51053),
    (3, "ROMERO MASÍAS M.",           "BC", "BC GRUPO 2", 16, 24.9,  35944),
    (4, "CALDERÓN CÓRDOVA K.",        "BEP","BEP ZONAL 3",  2, 54.5,  17595),
    (5, "DUHART NARANJO L. M.",       "BC", "BC GRUPO 1",  2, 20.7,  15656),
    (6, "CABRERA GALÁN DE NOLASCO Z.","BEL","BEL ZONAL 1",  2, 22.5,  11897),
    (7, "CHANG PILLIHUARI D. Z.",     "BEL","BEL ZONAL 2",  2, 49.5,   9535),
    (8, "BOLAÑOS DÍAZ R. F.",         "BEL","BEL ZONAL 2",  2, 49.0,   8570),
    (9, "MEDRANO LUNA A. A.",         "BEL","BEL ZONAL 1",  1, 19.7,   8348),
    (10,"MORA HARTMANN R. M.",        "BC", "BC GRUPO 2",  3, 42.5,   3911),
]

# Barras proporcionales
max_val = 1933255
bar_max_w = Inches(6.5)

rhdrs = ["#", "Ejecutivo", "Bca", "Clientes", "Cumpl.", "Extra Perdido"]
rcws  = [Inches(0.4), Inches(3.5), Inches(0.6), Inches(0.8), Inches(0.8), Inches(1.4)]
rxs   = [Inches(0.3)]
for cw2 in rcws[:-1]:
    rxs.append(rxs[-1]+cw2+Inches(0.05))
bar_start = rxs[-1]+rcws[-1]+Inches(0.2)

for j, (hdr, cw2, x2) in enumerate(zip(rhdrs, rcws, rxs)):
    rect(sl7, x2, Inches(1.4), cw2, Inches(0.35), fill=AZUL_OSC)
    txbox(sl7, hdr, x2+Inches(0.03), Inches(1.42), cw2-Inches(0.06), Inches(0.32),
          size=10, bold=True, color=BLANCO)
txbox(sl7, "Barra comparativa", bar_start, Inches(1.42), Inches(2.0), Inches(0.32),
      size=10, bold=True, color=AZUL_OSC)

banca_colors = {"BC": AZUL_OSC, "BEL": AZUL_MED, "BEP": VERDE_OSC}

for i, (rank, ejec, bca, zonal, cli, cumpl, extra) in enumerate(all_ejecs):
    ry = Inches(1.78) + i*Inches(0.47)
    bg = RGBColor(0xF0, 0xF4, 0xFF) if i%2==0 else BLANCO
    for cw2, x2 in zip(rcws, rxs):
        rect(sl7, x2, ry, cw2, Inches(0.44), fill=bg)
    vals = [str(rank), ejec, bca, str(cli), f"{cumpl:.1f}%", f"S/ {extra:,.0f}"]
    for j, (v, cw2, x2) in enumerate(zip(vals, rcws, rxs)):
        vc = ROJO_ALERT if j==5 else (banca_colors.get(v, NEGRO) if j==2 else NEGRO)
        bold = j==5 or j==0
        txbox(sl7, v, x2+Inches(0.03), ry+Inches(0.05), cw2-Inches(0.06), Inches(0.35),
              size=11, color=vc, bold=bold)
    # Barra
    bar_w = bar_max_w * (extra / max_val)
    if bar_w < Inches(0.05): bar_w = Inches(0.05)
    bar_color = banca_colors.get(bca, AZUL_OSC)
    rect(sl7, bar_start, ry+Inches(0.1), bar_w, Inches(0.25), fill=bar_color)

# Leyenda
legend_items = [("BC", AZUL_OSC), ("BEL", AZUL_MED), ("BEP", VERDE_OSC)]
lx = bar_start
for lbl, col in legend_items:
    rect(sl7, lx, Inches(6.55), Inches(0.2), Inches(0.2), fill=col)
    txbox(sl7, lbl, lx+Inches(0.25), Inches(6.55), Inches(0.6), Inches(0.25),
          size=10, color=NEGRO)
    lx += Inches(0.9)

# ── SLIDE 8: CIERRE ──────────────────────────────────────────────────────────
sl8 = add_slide()
rect(sl8, 0, 0, W, H, fill=AZUL_OSC)

circ3 = sl8.shapes.add_shape(9, Inches(-1), Inches(-1), Inches(5), Inches(5))
circ3.fill.solid(); circ3.fill.fore_color.rgb = AZUL_MED
circ3.line.fill.background()
circ4 = sl8.shapes.add_shape(9, Inches(10), Inches(4.5), Inches(5), Inches(5))
circ4.fill.solid(); circ4.fill.fore_color.rgb = VERDE_OSC
circ4.line.fill.background()

txbox(sl8, "PRÓXIMOS PASOS", Inches(1.5), Inches(1.2), Inches(10), Inches(0.6),
      size=14, bold=True, color=VERDE_PRI, align=PP_ALIGN.CENTER)
txbox(sl8, "Acción sobre incumplimientos", Inches(1.5), Inches(1.8), Inches(10), Inches(1.0),
      size=36, bold=True, color=BLANCO, align=PP_ALIGN.CENTER)

steps = [
    ("01", "Priorizar BC GRUPO 2", "Concentra S/1.97M de la pérdida. Revisión urgente con Ejecutivo Rodríguez."),
    ("02", "Seguimiento semanal", "Establecer checkpoint semanal por zonal hasta levantar cumplimiento."),
    ("03", "Plan de acción por ejecutivo", "Cada ejecutivo propone acciones concretas con fecha de compromiso."),
]
sw = Inches(3.6)
sx = Inches(0.8)
for i, (num, titulo, desc) in enumerate(steps):
    cx2 = sx + i*(sw+Inches(0.25))
    rect(sl8, cx2, Inches(3.0), sw, Inches(3.0), fill=AZUL_MED)
    txbox(sl8, num, cx2, Inches(3.05), sw, Inches(0.7),
          size=36, bold=True, color=VERDE_PRI, align=PP_ALIGN.CENTER)
    txbox(sl8, titulo, cx2+Inches(0.15), Inches(3.7), sw-Inches(0.3), Inches(0.55),
          size=14, bold=True, color=BLANCO, align=PP_ALIGN.CENTER)
    txbox(sl8, desc, cx2+Inches(0.15), Inches(4.3), sw-Inches(0.3), Inches(1.5),
          size=12, color=RGBColor(0xCC, 0xDD, 0xFF), align=PP_ALIGN.CENTER)

txbox(sl8, "Interbank · Product Management · Junio 2026",
      Inches(0), Inches(6.85), W, Inches(0.4),
      size=10, color=RGBColor(0x88, 0x99, 0xBB), align=PP_ALIGN.CENTER)

# ── GUARDAR ───────────────────────────────────────────────────────────────────
out = r"C:\00.OrdenTotal\18.SecondBrain\PARA\areas\product-management\docs\seguimiento-compromisos-depositos-abr2026.pptx"
prs.save(out)
print("Guardado:", out)
